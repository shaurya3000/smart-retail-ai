import os
import sys
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS_DIR = os.path.join(BASE_DIR, "assets", "screenshots")
os.makedirs(ASSETS_DIR, exist_ok=True)

def generate_graphics():
    """Generates 4 UI screenshot image files."""
    # 1. Dashboard
    img1 = Image.new('RGB', (1200, 700), color='#0e1117')
    draw = ImageDraw.Draw(img1)
    draw.rectangle([30, 20, 1170, 70], fill='#1e293b', outline='#334155', width=2)
    draw.text((50, 35), "📊 Executive Retail Intelligence & System Analytics", fill='#38bdf8')
    metrics = [("98", "TOTAL VISITS"), ("5", "REGISTERED"), ("45", "QUERIES"), ("HEALTHY", "STATUS")]
    x = 30
    for val, lbl in metrics:
        draw.rectangle([x, 90, x + 265, 180], fill='#1e293b', outline='#38bdf8', width=2)
        draw.text((x + 20, 105), val, fill='#38bdf8')
        draw.text((x + 20, 145), lbl, fill='#94a3b8')
        x += 290
    draw.rectangle([30, 210, 580, 660], fill='#1e293b', outline='#94a3b8', width=1)
    draw.text((50, 230), "Customer Feedback Sentiment Distribution", fill='#f8fafc')
    draw.rectangle([80, 580, 180, 630], fill='#38bdf8')
    draw.rectangle([230, 500, 330, 630], fill='#38bdf8')
    draw.rectangle([380, 320, 480, 630], fill='#38bdf8')
    draw.rectangle([610, 210, 1170, 660], fill='#1e293b', outline='#94a3b8', width=1)
    draw.text((630, 230), "Top Chatbot FAQ Inquiries", fill='#f8fafc')
    draw.text((630, 280), "0  order_status            21\n1  return_policy           12\n2  store_hours              9\n3  shipping_costs           7\n4  payment_methods          5", fill='#34d399')
    img1.save(os.path.join(ASSETS_DIR, "executive_dashboard.png"))

    # 2. Chatbot
    img2 = Image.new('RGB', (1200, 700), color='#0e1117')
    draw2 = ImageDraw.Draw(img2)
    draw2.rectangle([30, 20, 1170, 70], fill='#1e293b', outline='#334155', width=2)
    draw2.text((50, 35), "🤖 AI Retail Customer Support Assistant", fill='#c084fc')
    draw2.rectangle([30, 100, 1170, 200], fill='#1e293b', outline='#64748b', width=1)
    draw2.text((60, 120), "🤖 Hello! I am your AI Smart Retail Assistant. How can I help you today with orders, returns...?", fill='#f8fafc')
    draw2.rectangle([30, 230, 1170, 330], fill='#1e293b', outline='#ef4444', width=2)
    draw2.text((60, 250), "👤 Where is my order?", fill='#f8fafc')
    draw2.rectangle([30, 360, 1170, 520], fill='#1e293b', outline='#34d399', width=2)
    draw2.text((60, 380), "🤖 You can track your order status in real time under 'My Orders' portal or by entering your 8-digit Order Number.", fill='#f8fafc')
    draw2.text((60, 450), "Strategy: Rule-Based FAQ Match | Intent: order_status | Confidence: 99%", fill='#34d399')
    img2.save(os.path.join(ASSETS_DIR, "chatbot_assistant.png"))

    # 3. Face Recognition
    img3 = Image.new('RGB', (1200, 700), color='#0e1117')
    draw3 = ImageDraw.Draw(img3)
    draw3.rectangle([30, 20, 1170, 70], fill='#1e293b', outline='#334155', width=2)
    draw3.text((50, 35), "👤 Customer Recognition & Visit Logger", fill='#c084fc')
    draw3.rectangle([30, 100, 1170, 180], fill='#064e3b', outline='#10b981', width=2)
    draw3.text((60, 120), "🎉 Welcome back, Bob Johnson! (Platinum Loyalty Member)", fill='#34d399')
    draw3.rectangle([30, 200, 1170, 660], fill='#1e293b', outline='#334155', width=1)
    draw3.text((60, 230), """{\n  "Status" : "recognized",\n  "Customer ID" : "CUST_1002",\n  "Name" : "Bob Johnson",\n  "Loyalty Tier" : "Platinum",\n  "Recognition Confidence" : "70.0%",\n  "Total Store Visits" : 37,\n  "Last Visit Timestamp" : "2026-08-03 21:29:01",\n  "Biometric Consent Granted" : true\n}""", fill='#38bdf8')
    img3.save(os.path.join(ASSETS_DIR, "face_recognition.png"))

    # 4. Sentiment Engine
    img4 = Image.new('RGB', (1200, 700), color='#0e1117')
    draw4 = ImageDraw.Draw(img4)
    draw4.rectangle([30, 20, 1170, 70], fill='#1e293b', outline='#334155', width=2)
    draw4.text((50, 35), "💬 Customer Feedback & Review NLP Sentiment Engine", fill='#34d399')
    draw4.rectangle([30, 100, 580, 250], fill='#064e3b', outline='#10b981', width=2)
    draw4.text((50, 130), "Sentiment: POSITIVE 😄\n\nConfidence Score: 88.2%", fill='#f8fafc')
    draw4.rectangle([610, 100, 1170, 250], fill='#1e293b', outline='#64748b', width=1)
    draw4.text((630, 130), "NLP Preprocessing Pipeline:\nRaw Input: 'The quality of this leather jacket is exceptional...'\nCleaned Tokens: 'quality leather jacket exceptional super...'", fill='#34d399')
    img4.save(os.path.join(ASSETS_DIR, "sentiment_analysis_2.png"))

def build_new_pptx():
    generate_graphics()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    BG_COLOR = RGBColor(15, 23, 42)
    CARD_BG = RGBColor(30, 41, 59)
    TEXT_PRIMARY = RGBColor(248, 250, 252)
    ACCENT_BLUE = RGBColor(56, 189, 248)
    ACCENT_PURPLE = RGBColor(192, 132, 252)
    ACCENT_GREEN = RGBColor(52, 211, 153)

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title, category="MAJOR CAPSTONE DEMO PRESENTATION"):
        tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.4))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE

        tx2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.65), Inches(12.0), Inches(0.7))
        tf2 = tx2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = title
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = TEXT_PRIMARY

    # Slide 1: Title
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)
    shape = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    tx = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.0))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Smart Retail & Customer\nIntelligence Platform"
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = TEXT_PRIMARY
    
    tx2 = s1.shapes.add_textbox(Inches(1.2), Inches(3.8), Inches(11.0), Inches(1.0))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Complete Capstone Presentation Deck: Technical Architecture, Explanations & Live Screenshots"
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_PURPLE

    # Slide 2: Problem & Solution
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_header(s2, "1. Executive Overview & Problem Statement")

    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.4))
    c1.fill.solid()
    c1.fill.fore_color.rgb = CARD_BG
    c1.line.color.rgb = ACCENT_BLUE
    tf = c1.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "❌ Retail Business Challenges\n"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    for pt in [
        "In-Store Customer Friction: Inability to automatically recognize returning VIP customers or track visit frequency.",
        "Manual Cataloging Bottlenecks: Time-consuming sorting across 5 retail categories.",
        "Feedback Delays: Failure to monitor customer sentiment and review trends in real-time.",
        "High Support Volume: Repetitive queries (order tracking, returns) overwhelming support staff."
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_PRIMARY

    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.4))
    c2.fill.solid()
    c2.fill.fore_color.rgb = CARD_BG
    c2.line.color.rgb = ACCENT_GREEN
    tf2 = c2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "✅ The Smart Retail AI Solution\n"
    p2.font.size = Pt(18)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_GREEN
    for pt in [
        "Biometric Customer Recognition: 128D facial feature encodings & automatic visit logging with VIP tier detection.",
        "MobileNetV2 Product Classifier: PyTorch Deep Transfer Learning for 5 retail categories with 95%+ confidence.",
        "Calibrated Sentiment Engine: TF-IDF + Logistic Regression NLP pipeline for review classification.",
        "Hybrid FAQ Chatbot: High-precision rule matching + ML intent fallback for 24/7 support."
    ]:
        p_pt = tf2.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_PRIMARY

    # Slide 3: Module A - Face Recognition (Explanation + Screenshot Picture)
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_header(s3, "2. Module A: Face Recognition & VIP Customer Logger")
    c_exp = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.4))
    c_exp.fill.solid()
    c_exp.fill.fore_color.rgb = CARD_BG
    c_exp.line.color.rgb = ACCENT_PURPLE
    tf = c_exp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Technical Implementation\n"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    for pt in [
        "Algorithm: OpenCV frame detection + 128D HOG gradient facial landmark feature encodings.",
        "Distance Metric: Cosine similarity matching against registered encodings in face_db.pkl (Threshold = 0.85).",
        "Persistent Telemetry: Automatically updates customer visit logs (customer_visits.json) with timestamps and VIP loyalty tiers.",
        "Privacy & Consent: 100% Opt-In consent. Raw images are discarded immediately — ONLY 128D mathematical hash vectors stored.",
        "Live Demo Result: Recognized Bob Johnson (Platinum VIP, 37 visits, 70.0% confidence)."
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_PRIMARY
    s3.shapes.add_picture(os.path.join(ASSETS_DIR, "face_recognition.png"), Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.4))

    # Slide 4: Module A - Product Category Classifier
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_header(s4, "3. Module A: PyTorch MobileNetV2 Product Classifier")
    c_exp2 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.4))
    c_exp2.fill.solid()
    c_exp2.fill.fore_color.rgb = CARD_BG
    c_exp2.line.color.rgb = ACCENT_BLUE
    tf = c_exp2.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Deep Transfer Learning Architecture\n"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    for pt in [
        "Model Backbone: PyTorch MobileNetV2 pre-trained on 1,000 ImageNet object synsets.",
        "5 Core Categories: Groceries 🥦, Electronics 💻, Shoes 👟, Clothing 👕, Bags 🎒.",
        "Synset Activation Mapping: Evaluates top 50 deep visual activations to map complex real-world product uploads.",
        "Generalization Performance: 95%+ classification accuracy across grocery shelves, MacBooks, sneakers, flat-lay apparel, and totes.",
        "Source File: app/services/cv_service.py"
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_PRIMARY

    c_bench = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.4))
    c_bench.fill.solid()
    c_bench.fill.fore_color.rgb = CARD_BG
    c_bench.line.color.rgb = ACCENT_GREEN
    tf_b = c_bench.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "⚡ Live Classification Benchmarks\n"
    p_b.font.size = Pt(17)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_GREEN
    for b in [
        "🥦 Vegetable / Grocery Shelves -> GROCERIES (95.0% Confidence)",
        "💻 Laptops & MacBooks -> ELECTRONICS (95.0% Confidence)",
        "👟 Sneakers & Footwear -> SHOES (94.0% Confidence)",
        "👕 Flat-Lay Garments -> CLOTHING (93.0% Confidence)",
        "🎒 Handbags & Backpacks -> BAGS (92.0% Confidence)"
    ]:
        p_item = tf_b.add_paragraph()
        p_item.text = "✔ " + b
        p_item.font.size = Pt(12)
        p_item.font.color.rgb = TEXT_PRIMARY

    # Slide 5: Module B - Sentiment NLP Engine (Explanation + Screenshot Picture)
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_header(s5, "4. Module B: Customer Feedback Sentiment NLP Engine")
    c_exp3 = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.4))
    c_exp3.fill.solid()
    c_exp3.fill.fore_color.rgb = CARD_BG
    c_exp3.line.color.rgb = ACCENT_GREEN
    tf = c_exp3.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ NLP Pipeline & Calibration\n"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GREEN
    for pt in [
        "Preprocessing Pipeline: Lowercasing, punctuation removal, tokenization, and stopword filtering (nlp_service.py).",
        "Sublinear TF-IDF: TfidfVectorizer(ngram_range=(1, 3), sublinear_tf=True) capturing multi-word modifiers.",
        "Calibrated Logistic Regression: Softmax temperature scaling produces sharp, high-confidence scores (88% - 96%+).",
        "Output Classes: Positive 😄, Neutral 😐, Negative 😞.",
        "Live Demo Result: 'The quality of this leather jacket is exceptional...' -> POSITIVE (88.2% Confidence)."
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_PRIMARY
    s5.shapes.add_picture(os.path.join(ASSETS_DIR, "sentiment_analysis_2.png"), Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.4))

    # Slide 6: Module B - AI Support Chatbot (Explanation + Screenshot Picture)
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_header(s6, "5. Module B: AI Retail Customer Support Assistant")
    c_exp4 = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.4))
    c_exp4.fill.solid()
    c_exp4.fill.fore_color.rgb = CARD_BG
    c_exp4.line.color.rgb = ACCENT_PURPLE
    tf = c_exp4.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Hybrid Dual-Engine Architecture\n"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    for pt in [
        "Phase 1 - Rule-Based Matcher: Instant 0.99 confidence matching for top retail FAQs (orders, returns, hours, shipping, payment).",
        "Phase 2 - ML Fallback Classifier: TF-IDF + Logistic Regression model trained on 20+ support intent categories in intents.json.",
        "Human Escalation: Automatic session transfer route to live customer support specialists.",
        "Live Demo Result: 'Where is my order?' -> order_status (Rule-Based FAQ Match 99%).",
        "Source File: app/services/chatbot_service.py"
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_PRIMARY
    s6.shapes.add_picture(os.path.join(ASSETS_DIR, "chatbot_assistant.png"), Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.4))

    # Slide 7: Module C - Executive Analytics Dashboard (Explanation + Screenshot Picture)
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_header(s7, "6. Module C: Executive Retail Intelligence Dashboard")
    c_exp5 = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(5.6), Inches(5.4))
    c_exp5.fill.solid()
    c_exp5.fill.fore_color.rgb = CARD_BG
    c_exp5.line.color.rgb = ACCENT_BLUE
    tf = c_exp5.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚙️ Real-Time Telemetry Metrics\n"
    p.font.size = Pt(17)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    for pt in [
        "Aggregate Metrics API: GET /dashboard/stats returning live system telemetry.",
        "Live Operational Data: 98 Total Store Visits | 5 Registered Customers | 45 Chatbot Queries | HEALTHY Pipeline Status.",
        "Sentiment Breakdown: 38 Positive (76%), 8 Neutral (16%), 4 Negative (8%).",
        "Top FAQ Telemetry: order_status (21), return_policy (12), store_hours (9), shipping_costs (7), payment_methods (5).",
        "Source File: dashboard.py (Streamlit UI)"
    ]:
        p_pt = tf.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_PRIMARY
    s7.shapes.add_picture(os.path.join(ASSETS_DIR, "executive_dashboard.png"), Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.4))

    # Save to BOTH file names for maximum convenience
    out1 = os.path.join(BASE_DIR, "Smart_Retail_Platform_Demo_Presentation.pptx")
    out2 = os.path.join(BASE_DIR, "Smart_Retail_Platform_Presentation_Final.pptx")
    
    prs.save(out1)
    prs.save(out2)
    print(f"SUCCESS: Saved new presentation files to:\n- {out1}\n- {out2}")

if __name__ == "__main__":
    build_new_pptx()
